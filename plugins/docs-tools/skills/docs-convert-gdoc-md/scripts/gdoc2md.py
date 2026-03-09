"""
Export Google Docs to Markdown, Slides to Markdown (via PPTX), or Sheets to CSV.

Requires gcloud CLI and python-pptx (for Slides export).

python gdoc2md.py <google-doc-or-slides-or-sheets-url> [output]
"""

import re
import subprocess
import sys
from io import BytesIO
from pathlib import Path
from urllib.request import Request, urlopen
from urllib.error import HTTPError


def check_dependencies():
    result = subprocess.run(
        ["gcloud", "version"], capture_output=True,
    )
    if result.returncode != 0:
        print("Error: gcloud CLI is not installed.", file=sys.stderr)
        install = "https://cloud.google.com/sdk/docs/install"
        print(f"  Install: {install}", file=sys.stderr)
        sys.exit(1)


def detect_mode(url):
    if "/presentation/d/" in url:
        return "slides"
    if "/spreadsheets/d/" in url:
        return "sheets"
    return "doc"


def parse_args():
    if len(sys.argv) < 2:
        print(
            f"Usage: {sys.argv[0]} "
            "<google-doc-or-slides-or-sheets-url> [output]"
        )
        sys.exit(1)

    url = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else ""
    mode = detect_mode(url)
    return url, output, mode


def extract_id(url, output, mode):
    match = re.search(r"/d/([a-zA-Z0-9_-]+)", url)
    if not match:
        print(f"Error: Could not extract ID from URL: {url}", file=sys.stderr)
        sys.exit(1)

    file_id = match.group(1)
    if not output:
        extensions = {"slides": ".md", "sheets": ".csv", "doc": ".md"}
        output = f"{file_id}{extensions[mode]}"

    return file_id, output


def authenticate():
    result = subprocess.run(
        ["gcloud", "auth", "print-access-token"],
        capture_output=True, text=True,
    )
    if result.returncode == 0:
        return
    print("Authenticating with Google...")
    subprocess.run(
        ["gcloud", "auth", "login", "--enable-gdrive-access"],
        check=True,
    )


def get_token():
    return subprocess.run(
        ["gcloud", "auth", "print-access-token"],
        capture_output=True, text=True, check=True,
    ).stdout.strip()


def download(url, token):
    req = Request(url, headers={"Authorization": f"Bearer {token}"})
    try:
        with urlopen(req) as resp:
            return resp.read()
    except HTTPError as e:
        messages = {
            401: "Authentication failed (401). "
                 "Try: gcloud auth login "
                 "--enable-gdrive-access",
            403: "Access denied (403). "
                 "Check you have permission.",
            404: "Not found (404). "
                 "Check the URL is correct.",
        }
        msg = messages.get(e.code, f"HTTP {e.code}")
        print(f"Error: {msg}", file=sys.stderr)
        sys.exit(1)


def pptx_to_markdown(data):
    """Convert PPTX bytes to structured Markdown."""
    try:
        from pptx import Presentation
    except ImportError:
        print(
            "Error: python-pptx is required for Slides export.",
            file=sys.stderr,
        )
        print("  Install: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(BytesIO(data))
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        lines.append("")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    level = paragraph.level
                    if level > 0:
                        indent = "  " * (level - 1)
                        lines.append(f"{indent}- {text}")
                    else:
                        lines.append(text)
                lines.append("")

            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    cells = [
                        cell.text.strip().replace("|", "\\|")
                        for cell in row.cells
                    ]
                    lines.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        lines.append(
                            "| " + " | ".join(["---"] * len(cells)) + " |"
                        )
                lines.append("")

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("> **Notes:** " + notes_text)
                lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


def fetch(file_id, output, mode):
    token = get_token()
    base = "https://docs.google.com"

    if mode == "slides":
        export_url = (
            f"{base}/presentation/d/{file_id}/export?format=pptx"
        )
    elif mode == "sheets":
        export_url = (
            f"{base}/spreadsheets/d/{file_id}/export?format=csv"
        )
    else:
        export_url = (
            f"{base}/document/d/{file_id}/export?format=md"
        )

    data = download(export_url, token)
    output_path = Path(output)

    if mode == "slides":
        markdown = pptx_to_markdown(data)
        output_path.write_text(markdown)
    else:
        output_path.write_bytes(data)

    print(f"Saved to {output}")


def main():
    url, output, mode = parse_args()
    check_dependencies()
    file_id, output = extract_id(url, output, mode)
    authenticate()
    fetch(file_id, output, mode)


if __name__ == "__main__":
    main()
